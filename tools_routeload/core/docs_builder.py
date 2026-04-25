"""生成程序使用说明书 (DOCX) 和 PPT (PPTX) 资源文件。
首次访问时按需生成；每次启动会重新生成保证内容最新。
"""
from __future__ import annotations

from io import BytesIO

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PRGBColor
from pptx.util import Inches, Pt as PPt

# ─────────────────────────────────────────
# 共用：功能介绍 / 计算逻辑文本
# ─────────────────────────────────────────

INTRO_BULLETS = [
    "批量解析 OFP 飞行计划 TXT，一次可处理 500+ 份",
    "自动按航线（起飞-目的-线路）分组，同航线下按机型 A319 → A320-214W → A320-251 顺序汇总",
    "自动生成 Word 报告（纵向 A4，含大标题 + 副标题 + 完整表格 + 边框）",
    "机场四字码与机型代号映射可在 config 目录的 JSON 文件里随时增删",
    "支持网页版（任何人浏览器打开即用）与本地版（双击启动）",
]

LOGIC_ROWS = [
    ("大标题", "文件名机场四字码 + 末位字母→线路", "S=南线 / N=北线 / W=W线"),
    ("月份", "文件名末两位数字", "ZSWX-ZWTL S07 → 7"),
    ("机型副标题", "TXT 第 3 行机号代号 → 映射表", "B306C → A319-115"),
    ("起飞重量", "TXT 中 TOW 后数字", "70000"),
    ("总加油量", "TXT 中 TOTL 后数字", "17700"),
    ("航程油量", "DEST 行 FUEL 列", "12283"),
    ("航程时间", "DEST 行 TIME 列", "04/40"),
    ("航线距离", "DEST 行 DIST 列", "1818"),
    ("最大业载", "TXT 中 AV PLD 后数字", "11118"),
    ("航路平均风", "ROUTE AVG WIND 后字段", "M038"),
    ("额外油", "TXT 中 XTRA 后数字", "768"),
    ("落地剩油", "TARGET ARRIVAL 后数字", "5000"),
    ("计算高度", "FLIGHT LEVEL 字段后续 4 行内最大 FL × 100", "FL301 → 30100 FT"),
    ("限重计算温度", "固定值", "0"),
    ("人数", "AV PLD ÷ 85 向下取整", "11118 / 85 = 130"),
]


# ─────────────────────────────────────────
# DOCX 说明书
# ─────────────────────────────────────────

def _h(doc: Document, text: str, size: int, bold: bool = True, color=None):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(6)
    p.paragraph_format.space_after = Pt(4)
    run = p.add_run(text)
    run.font.name = "宋体"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    run.font.size = Pt(size)
    run.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _para(doc: Document, text: str, size: int = 10):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(2)
    run = p.add_run(text)
    run.font.name = "宋体"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    run.font.size = Pt(size)


def build_manual_bytes() -> bytes:
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "宋体"
    style.element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    style.font.size = Pt(10.5)

    _h(doc, "航线载量分析系统 · 使用说明书", 18, color=RGBColor(0x1A, 0x52, 0x76))
    _para(doc, "湖南航空公司 · 王迪 出品 · 技术支持：杨清云")

    _h(doc, "一、功能介绍", 14, color=RGBColor(0x1A, 0x52, 0x76))
    for s in INTRO_BULLETS:
        _para(doc, "• " + s)

    _h(doc, "二、操作流程", 14, color=RGBColor(0x1A, 0x52, 0x76))
    for i, step in enumerate([
        "打开网页 https://wangdi-payload.streamlit.app  或本地双击「启动.command」",
        "在「上传 TXT」区域拖入或选择 OFP TXT 文件（支持一次性 500+ 份批量上传）",
        "等待解析完成，预览数据表，确认无误",
        "点击「生成报告」按钮，再点「下载 Word」即可导出",
    ], 1):
        _para(doc, f"{i}. {step}")

    _h(doc, "三、字段抓取与计算逻辑", 14, color=RGBColor(0x1A, 0x52, 0x76))
    table = doc.add_table(rows=1 + len(LOGIC_ROWS), cols=3)
    hdr = table.rows[0].cells
    for i, t in enumerate(["输出字段", "数据来源 / 算法", "示例"]):
        hdr[i].text = t
        for r in hdr[i].paragraphs[0].runs:
            r.bold = True
            r.font.size = Pt(10)
    for r, (a, b, c) in enumerate(LOGIC_ROWS, 1):
        cells = table.rows[r].cells
        cells[0].text = a
        cells[1].text = b
        cells[2].text = c
        for cell in cells:
            for p in cell.paragraphs:
                for run in p.runs:
                    run.font.name = "宋体"
                    run._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
                    run.font.size = Pt(9.5)
    table.style = "Light Grid Accent 1"

    _h(doc, "四、文件命名规则", 14, color=RGBColor(0x1A, 0x52, 0x76))
    _para(doc, "示例：306C ZSWX-ZWTL S07.txt")
    _para(doc, "• 第一段：机号（如 306C）")
    _para(doc, "• 第二段：起飞机场-目的机场（ICAO 四字码）")
    _para(doc, "• 第三段（可选）：线路标识 + 月份，例如 S07 表示南线 7 月")

    _h(doc, "五、输出说明", 14, color=RGBColor(0x1A, 0x52, 0x76))
    _para(doc, "• Word 文档为纵向 A4 排版")
    _para(doc, "• 同一航线（起飞-目的-线路相同）汇总在同一个大标题下")
    _para(doc, "• 大标题下按机型 A319-115 → A320-214W → A320-251 顺序排列各机型表格")

    _h(doc, "六、自定义机场 / 机型映射", 14, color=RGBColor(0x1A, 0x52, 0x76))
    _para(doc, "如果遇到未录入的机场代号或机型代号，请编辑：")
    _para(doc, "• config/airports.json — ICAO 四字码到中文名映射")
    _para(doc, "• config/aircraft.json — 机号代号到机型名映射")

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────
# PPTX 介绍幻灯片 — 湖南航空风格（红色主题，16:9）
# ─────────────────────────────────────────

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# 湖南航空品牌色
HNA_RED = PRGBColor(0xC0, 0x00, 0x00)        # 主红
HNA_RED_LIGHT = PRGBColor(0xD7, 0x00, 0x3A)  # 标题红
HNA_GRAY = PRGBColor(0x92, 0x9A, 0x9F)       # 装饰灰
TEXT_DARK = PRGBColor(0x51, 0x54, 0x57)      # 正文深灰
TEXT_MUTED = PRGBColor(0x80, 0x80, 0x80)
TEXT_LIGHT = PRGBColor(0x96, 0x96, 0x96)
WHITE = PRGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FOOTER_SLOGAN = "团结  奋斗  高效  务实"


def _add_textbox(slide, left, top, width, height, text, *,
                 size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT,
                 font_name="微软雅黑"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = PPt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def _add_rect(slide, left, top, width, height, fill_color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if not line:
        shp.line.fill.background()
    return shp


def _add_footer(slide, page_no, total):
    # 底部细红线
    _add_rect(slide, Inches(0.5), Inches(7.05), Inches(12.33), Inches(0.03), HNA_RED)
    _add_textbox(slide, Inches(0.5), Inches(7.12), Inches(6), Inches(0.35),
                 FOOTER_SLOGAN, size=13, bold=True, color=TEXT_MUTED)
    _add_textbox(slide, Inches(11.5), Inches(7.12), Inches(1.5), Inches(0.35),
                 f"{page_no} / {total}", size=12, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)


def _add_section_header(slide, num, title):
    """每页顶部统一的章节标题条。"""
    # 红色竖条（加粗加高）
    _add_rect(slide, Inches(0.5), Inches(0.45), Inches(0.14), Inches(0.7), HNA_RED)
    _add_textbox(slide, Inches(0.78), Inches(0.4), Inches(12), Inches(0.85),
                 f"{num}  {title}", size=32, bold=True, color=HNA_RED_LIGHT)
    # 分隔线
    _add_rect(slide, Inches(0.5), Inches(1.25), Inches(12.33), Inches(0.02), HNA_GRAY)


def _new_slide(prs):
    blank = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(blank)


# ─────────── 幻灯片构造器 ────────────

def _slide_cover(prs, total):
    s = _new_slide(prs)
    # 右侧大红块
    _add_rect(s, Inches(8.66), Inches(0), Inches(4.67), Inches(7.5), HNA_RED)
    # 右上角小灰块装饰
    _add_rect(s, Inches(11.0), Inches(0), Inches(2.36), Inches(3.15), HNA_GRAY)
    # 顶部红色短线
    _add_rect(s, Inches(0.8), Inches(1.0), Inches(1.4), Inches(0.08), HNA_RED)
    # 主标题（超大）
    _add_textbox(s, Inches(0.8), Inches(1.3), Inches(7.5), Inches(2.4),
                 "航线载量分析系统", size=66, bold=True, color=HNA_RED_LIGHT)
    # 英文副标题
    _add_textbox(s, Inches(0.8), Inches(3.55), Inches(7.5), Inches(0.6),
                 "Route Payload Analysis System", size=22, color=TEXT_DARK)
    # 描述
    _add_textbox(s, Inches(0.8), Inches(4.4), Inches(7.5), Inches(1.5),
                 "OFP 飞行计划批量解析 · 自动汇总航线载量\n湖南航空 · 运行指挥部 · 签派室",
                 size=22, bold=True, color=TEXT_DARK)
    # 设计开发
    _add_textbox(s, Inches(0.8), Inches(5.95), Inches(5), Inches(0.55),
                 "系统开发：王迪", size=26, bold=True, color=TEXT_DARK)
    _add_textbox(s, Inches(0.8), Inches(6.5), Inches(5), Inches(0.5),
                 "技术支持：杨清云", size=20, color=TEXT_DARK)
    # 日期放右下白字
    _add_textbox(s, Inches(9), Inches(6.6), Inches(3.5), Inches(0.5),
                 "2026 年 4 月", size=18, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    # 页码
    _add_textbox(s, Inches(12.2), Inches(7.15), Inches(1), Inches(0.3),
                 f"1 / {total}", size=11, color=WHITE, align=PP_ALIGN.RIGHT)


def _slide_toc(prs, total, items):
    s = _new_slide(prs)
    _add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(1.2),
                 "目录", size=60, bold=True, color=HNA_RED_LIGHT)
    _add_textbox(s, Inches(0.5), Inches(1.65), Inches(12), Inches(0.5),
                 "CONTENTS", size=20, color=TEXT_MUTED)
    _add_rect(s, Inches(0.5), Inches(2.25), Inches(2.6), Inches(0.06), HNA_RED)
    # 两列布局，整张卡片更大更显眼
    cols = 2
    rows = (len(items) + cols - 1) // cols
    cell_w = 6.0
    cell_h = 4.5 / rows  # 撑满 2.6~7.0
    for i, item in enumerate(items):
        c, r = i % cols, i // cols
        x = 0.7 + c * cell_w
        y = 2.6 + r * cell_h
        # 编号（放大）
        _add_textbox(s, Inches(x), Inches(y), Inches(1.4), Inches(cell_h - 0.2),
                     f"{i+1:02d}", size=44, bold=True, color=HNA_RED)
        # 标题（放大）
        _add_textbox(s, Inches(x + 1.5), Inches(y + 0.15), Inches(cell_w - 1.6), Inches(cell_h - 0.4),
                     item, size=24, bold=True, color=TEXT_DARK)
    _add_footer(s, 2, total)


def _slide_section_bullets(prs, total, page_no, num, title, blocks):
    """blocks: list[(emoji, subtitle, [bullets])]"""
    s = _new_slide(prs)
    _add_section_header(s, num, title)
    n = len(blocks)
    if n <= 2:
        cols = n; rows = 1
    elif n <= 4:
        cols = 2; rows = 2
    else:
        cols = 3; rows = (n + 2) // 3
    # 卡片区：1.4 → 6.95 = 5.55in 高
    area_top = 1.45
    area_h = 5.5
    cell_w = 12.33 / cols
    cell_h = area_h / rows
    # 字号按行数自适应：行越多越紧凑
    if rows == 1:
        sz_emoji, sz_sub, sz_bul, line_h = 56, 30, 22, 0.55
    elif rows == 2:
        sz_emoji, sz_sub, sz_bul, line_h = 40, 22, 17, 0.42
    else:
        sz_emoji, sz_sub, sz_bul, line_h = 32, 18, 14, 0.34
    for idx, (emoji, sub, bullets) in enumerate(blocks):
        c = idx % cols; r = idx // cols
        x = 0.5 + c * cell_w
        y = area_top + r * cell_h
        # 卡片背景（占满，仅留 0.1 间隙）
        _add_rect(s, Inches(x + 0.08), Inches(y + 0.05), Inches(cell_w - 0.16), Inches(cell_h - 0.1),
                  PRGBColor(0xF5, 0xF5, 0xF5))
        # emoji（左上大）
        _add_textbox(s, Inches(x + 0.3), Inches(y + 0.18), Inches(1.4), Inches(sz_emoji * 0.025),
                     emoji, size=sz_emoji, bold=True, color=HNA_RED)
        # 副标题（右侧，与 emoji 同行）
        _add_textbox(s, Inches(x + 1.5), Inches(y + 0.3), Inches(cell_w - 1.7), Inches(0.7),
                     sub, size=sz_sub, bold=True, color=HNA_RED_LIGHT)
        # 红色短线
        line_y = y + 0.18 + sz_emoji * 0.025 + 0.15
        _add_rect(s, Inches(x + 0.3), Inches(line_y), Inches(0.8), Inches(0.04), HNA_RED)
        # bullets — 每条占 line_h 高度
        bullet_top = line_y + 0.18
        for bi, b in enumerate(bullets):
            _add_textbox(s, Inches(x + 0.3), Inches(bullet_top + bi * line_h),
                         Inches(cell_w - 0.6), Inches(line_h),
                         f"·  {b}", size=sz_bul, color=TEXT_DARK)
    _add_footer(s, page_no, total)


def _slide_steps(prs, total, page_no, num, title, steps):
    """steps: list[(no, name, desc)] — 横排流程图"""
    s = _new_slide(prs)
    _add_section_header(s, num, title)
    n = len(steps)
    cell_w = 12.33 / n
    circle_d = 1.6  # 加大圆
    for i, (no, name, desc) in enumerate(steps):
        x = 0.5 + i * cell_w
        # 圆形编号
        circle = s.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + cell_w/2 - circle_d/2), Inches(1.8),
            Inches(circle_d), Inches(circle_d),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = HNA_RED
        circle.line.fill.background()
        tf = circle.text_frame
        tf.margin_left = tf.margin_right = Inches(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(no)
        r.font.name = "微软雅黑"; r.font.size = PPt(48); r.font.bold = True; r.font.color.rgb = WHITE
        # 名称
        _add_textbox(s, Inches(x), Inches(3.65), Inches(cell_w), Inches(0.7),
                     name, size=26, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
        # 描述
        _add_textbox(s, Inches(x + 0.2), Inches(4.5), Inches(cell_w - 0.4), Inches(2.4),
                     desc, size=18, color=TEXT_DARK, align=PP_ALIGN.CENTER)
        # 箭头
        if i < n - 1:
            _add_textbox(s, Inches(x + cell_w - 0.4), Inches(2.2), Inches(0.8), Inches(0.8),
                         "→", size=44, bold=True, color=HNA_GRAY, align=PP_ALIGN.CENTER)
    _add_footer(s, page_no, total)


def _slide_table(prs, total, page_no, num, title, header, rows):
    s = _new_slide(prs)
    _add_section_header(s, num, title)
    n_cols = len(header)
    n_rows = len(rows) + 1
    left = Inches(0.5); top = Inches(1.45)
    width = Inches(12.33); height = Inches(5.55)
    tbl = s.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    # 表头
    for i, h in enumerate(header):
        cell = tbl.cell(0, i)
        cell.fill.solid(); cell.fill.fore_color.rgb = HNA_RED
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.name = "微软雅黑"; run.font.size = PPt(18); run.font.bold = True
        run.font.color.rgb = WHITE
    # 数据
    for r_idx, row in enumerate(rows, 1):
        bg = WHITE if r_idx % 2 else PRGBColor(0xF8, 0xF8, 0xF8)
        for i, val in enumerate(row):
            cell = tbl.cell(r_idx, i)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
            run = p.add_run(); run.text = str(val)
            run.font.name = "微软雅黑"; run.font.size = PPt(14); run.font.color.rgb = TEXT_DARK
    _add_footer(s, page_no, total)


def _slide_end(prs, total, page_no):
    s = _new_slide(prs)
    _add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), HNA_RED)
    _add_textbox(s, Inches(0.5), Inches(2.0), Inches(12.33), Inches(2.0),
                 "THANKS", size=140, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(s, Inches(0.5), Inches(4.5), Inches(12.33), Inches(0.8),
                 "感谢观看 · 敬请指正", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(s, Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.7),
                 "系统开发：王迪    技术支持：杨清云", size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(s, Inches(12.2), Inches(7.15), Inches(1), Inches(0.3),
                 f"{page_no} / {total}", size=11, color=WHITE, align=PP_ALIGN.RIGHT)


def build_ppt_bytes() -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    toc_items = [
        "项目背景与价值",
        "核心功能特性",
        "操作流程演示",
        "字段抓取逻辑",
        "文件命名规则",
        "自定义映射",
    ]
    total = 2 + len(toc_items) + 1  # 封面 + 目录 + 各章 + 致谢

    # 1. 封面
    _slide_cover(prs, total)
    # 2. 目录
    _slide_toc(prs, total, toc_items)

    # 3. 项目背景与价值
    _slide_section_bullets(prs, total, 3, "01", "项目背景与价值", [
        ("✈", "为什么需要这个系统",
         ["每月需重新汇总航线载量数据，工作量大",
          "OFP TXT 字段散布多处，人工提取易错",
          "同一航线多机型需分别整理，重复劳动",
          "Word 排版要求统一，手工调整费时"]),
        ("🎯", "本系统能做什么",
         ["一键拖入数百份 OFP TXT 文件",
          "自动按航线 + 机型分组汇总",
          "自动生成符合公司要求的 Word 报告",
          "从「半天工作」→「一分钟完成」"]),
    ])

    # 4. 核心功能
    _slide_section_bullets(prs, total, 4, "02", "核心功能特性", [
        ("📥", "批量解析",
         ["支持 500+ 份 TXT 同时解析",
          "解析速度秒级响应",
          "失败文件单独标记"]),
        ("🧭", "智能分组",
         ["按起飞-目的-线路自动归类",
          "南/北/W 线自动识别",
          "同航线汇总到一个标题"]),
        ("🛩", "机型分层",
         ["A319-115 → A320-214W → A320-251",
          "顺序固定，输出统一",
          "未识别机型可补充映射"]),
        ("📄", "Word 一键导出",
         ["纵向 A4，宋体",
          "标题 + 副标题 + 完整带边框表格",
          "下载即用，无需二次排版"]),
        ("🌐", "网页 + 本地共用",
         ["浏览器打开网页即用",
          "本地双击「启动.command」",
          "代码完全一致"]),
        ("🛠", "可扩展",
         ["机场词典 JSON 可编辑",
          "机型词典 JSON 可编辑",
          "无需重新打包"]),
    ])

    # 5. 操作流程
    _slide_steps(prs, total, 5, "03", "操作流程演示", [
        ("①", "打开系统", "网页：访问\nwangdi-payload\n.streamlit.app\n本地：双击「启动」"),
        ("②", "上传 TXT", "拖入或选择\nOFP 文件\n支持批量"),
        ("③", "核对数据", "查看预览表\n确认字段无误"),
        ("④", "生成报告", "点击「生成报告」\n自动汇总分组"),
        ("⑤", "下载使用", "下载 Word 文档\n打印或归档"),
    ])

    # 6. 字段抓取逻辑（表格）
    _slide_table(prs, total, 6, "04", "字段抓取与计算逻辑",
                 ["输出字段", "数据来源 / 算法", "示例"],
                 [list(r) for r in LOGIC_ROWS])

    # 7. 文件命名规则
    _slide_section_bullets(prs, total, 7, "05", "文件命名规则", [
        ("📝", "命名格式",
         ["示例：306C ZSWX-ZWTL S07.txt",
          "三段空格分隔",
          "中英文均可"]),
        ("①", "第一段：机号",
         ["如 306C / 302Y / 321U",
          "对应 aircraft.json 映射",
          "决定副标题机型"]),
        ("②", "第二段：航线",
         ["格式 ICAO-ICAO",
          "如 ZSWX-ZWTL",
          "对应 airports.json 映射"]),
        ("③", "第三段：线路+月份",
         ["S=南线 / N=北线 / W=W线",
          "末两位为月份（07 = 7 月）",
          "可省略，省略时无线路标识"]),
    ])

    # 8. 自定义映射
    _slide_section_bullets(prs, total, 8, "06", "自定义映射", [
        ("🛬", "机场词典",
         ["文件：config/airports.json",
          "格式：ICAO 四字码 → 中文名",
          '示例："ZWTL": "吐鲁番"',
          "未录入会原样显示，编辑即可"]),
        ("✈", "机型词典",
         ["文件：config/aircraft.json",
          "格式：机号代号 → 机型名",
          '示例："B306C": "A319-115"',
          "影响 Word 副标题与排序"]),
    ])

    # 9. 致谢
    _slide_end(prs, total, total)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
